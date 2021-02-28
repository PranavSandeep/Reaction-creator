__author__ = "Pranav Sandeep"  # I have no idea what I am doing.

import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from tkinter import *
from pptx.util import Inches
import threading
from threading import Thread
import random
import wikipedia
import os
from selenium import webdriver
import time
import urllib.request
from selenium.webdriver.common.keys import Keys

chrome_options = webdriver.ChromeOptions()
chrome_options.add_argument("--headless")
chrome_options.add_argument("--no-startup-window")

browser = webdriver.Chrome('Chrome Driver/chromedriver.exe',
                           chrome_options=chrome_options)

browser.get("https://www.google.com/")


class SampleApp(tk.Tk):
    def __init__(self):
        tk.Tk.__init__(self)
        self.prs = Presentation()
        style = ttk.Style(self)
        self.tk.call('source', 'azure dark.tcl')
        style.theme_use('azure')
        self.VarList = StringVar(self)
        self.VarList.set("Normal")
        self.SlideType = ["Normal", "Wikipedia"]
        self.SlideBox = tk.OptionMenu(self, self.VarList, *self.SlideType)
        self.Titlelabel = ttk.Label(self, text="What will the title be?")
        self.Subtitlelabel = ttk.Label(self,
                                   text="What will the subtitle be?(Not required if you are using wikipedia search or image search)")
        self.TopicLabel = ttk.Label(self,
                                text="What will the topic be?(not required if you are not using wikipedia search)")
        self.sentencesLabel = ttk.Label(self,
                                    text="How many sentances should the wikipedia summary be?(not required if you are not using wikipedia search)")
        self.PPTName = ttk.Label(self, text="PPT Name?")
        self.PPTNameEntry = ttk.Entry(self)
        self.Title = ttk.Entry(self)
        self.Subtitle = ttk.Entry(self)
        self.Topic = ttk.Entry(self)
        self.Sentences = ttk.Entry(self)
        self.GetImg = ttk.Button(self, text="Get images", command=self.Download_Images)
        self.button = ttk.Button(self, text="Add slide", command=self.CreatePPT)
        self.SlideBox.pack()
        self.Titlelabel.pack()
        self.Title.pack()
        self.Subtitlelabel.pack()
        self.Subtitle.pack()
        self.TopicLabel.pack()
        self.Topic.pack()
        self.sentencesLabel.pack()
        self.Sentences.pack()
        self.OpenButton = ttk.Button(self, text="OpenPPT", command=self.OpenPPT)
        self.GetImg.pack()
        self.button.pack()
        self.PPTName.pack()
        self.PPTNameEntry.pack()
        self.OpenButton.pack()

    def CreatePPT(self):

        if self.VarList.get() == "Normal":
            blank_slide_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(blank_slide_layout)

            title = slide.shapes.title

            subtitle = slide.placeholders[1]

            title.text = self.Title.get()
            subtitle.text = self.Subtitle.get()

            self.prs.save(self.PPTNameEntry.get() + ".pptx")


        elif self.VarList.get() == "Wikipedia":
            Fourth_layout = self.prs.slide_layouts[1]
            slide4 = self.prs.slides.add_slide(Fourth_layout)

            title4 = slide4.shapes.title
            Content4 = slide4.placeholders[1]

            title4.text = self.Title.get()
            Topic = self.Topic.get()
            Sentences = int(self.Sentences.get())
            try:
                Content4.text = wikipedia.summary(Topic, sentences=Sentences)

            except wikipedia.exceptions.DisambiguationError as e:
                print(
                    "Hmm... There are more than one topics about this topic. Perhaps try being more specific? Or this topic may "
                    "not exist")

            self.prs.save(self.PPTNameEntry.get() + ".pptx")

    def Download_Images(self):

        search = browser.find_element_by_name('q')

        search.send_keys(self.Topic.get(), Keys.ENTER)

        elem = browser.find_element_by_link_text('Images')
        elem.get_attribute('href')
        elem.click()

        value = 0
        for i in range(20):
            browser.execute_script("scrollBy('+ str(value) +',+1000);")
            value += 1000
            time.sleep(3)

        elem1 = browser.find_element_by_id('islmp')

        sub = elem1.find_elements_by_tag_name("img")

        try:
            os.mkdir('downloads')
        except FileExistsError:
            pass

        count = 0

        flag = 1

        for i in sub:

            if flag >= 2:
                break
            src = i.get_attribute('src')
            try:
                if src != None:
                    src = str(src)
                    print(src)
                    count += 1
                    urllib.request.urlretrieve(src, os.path.join('downloads', 'image' + "A" + '.jpg'))
                else:
                    raise TypeError
            except TypeError:
                print('fail')

            flag += 1

        Third_layout = self.prs.slide_layouts[6]
        slide3 = self.prs.slides.add_slide(Third_layout)

        left = top = Inches(1)

        try:
            pic = slide3.shapes.add_picture("downloads/" + "imageA" + ".jpg", left, top)

        except Exception as e:
            print(
                "Sorry, mate. But there is no wikipedia image for the topic you're searching for. Try inserting one manually.",
                e)

        self.prs.save(self.PPTNameEntry.get() + ".pptx")

    def OpenPPT(self):
        os.startfile(self.PPTNameEntry.get() + ".pptx")


app = SampleApp()

app.mainloop()
